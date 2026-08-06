import argparse
import os
import re
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Tuple

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from qdrant_client.http import models as qmodels
from sentence_transformers import SentenceTransformer

from embedding_utils import resolve_embedding_device, resolve_embedding_model
from graph_store import add_graph_store_args, create_graph_store_from_args
from doc_local_qdrant import get_document_qdrant_store
from project_contract import ProjectNotRegisteredError, resolve_project_targets

try:
    from dotenv import load_dotenv
except Exception:
    load_dotenv = None

from entity_extractors import (
    build_spacy_pipeline,
    extract_entities_gemini,
    extract_entities_gliner,
    extract_entities_gliner_batch,
    extract_entities_langextract,
    build_gliner_model,
    parse_gliner_labels,
)


def _load_env() -> None:
    if load_dotenv is None:
        return
    env_path = Path(__file__).with_name(".env")
    if env_path.exists():
        load_dotenv(dotenv_path=env_path)
    else:
        load_dotenv()


def read_pdf_text(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n\n".join(p for p in parts if p.strip())


def read_text_file(text_path: Path) -> str:
    return text_path.read_text(encoding="utf-8").strip()


def read_docx_text(docx_path: Path) -> str:
    doc = Document(str(docx_path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(parts)


def read_pptx_text(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                stripped = text.strip()
                if stripped:
                    parts.append(stripped)
    return "\n\n".join(parts)


def read_xlsx_text(xlsx_path: Path) -> str:
    workbook = load_workbook(filename=str(xlsx_path), data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = []
            for value in row:
                if value is None:
                    continue
                text = str(value).strip()
                if text:
                    cells.append(text)
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def split_paragraphs(text: str, max_chars: int) -> List[str]:
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n+", text) if p.strip()]
    chunks: List[str] = []
    for para in paragraphs:
        if len(para) <= max_chars:
            chunks.append(para)
        else:
            for i in range(0, len(para), max_chars):
                chunk = para[i : i + max_chars].strip()
                if chunk:
                    chunks.append(chunk)
    return chunks


def _set_langextract_overrides(model_id: str | None, model_url: str | None) -> None:
    if model_id:
        os.environ["LANGEXTRACT_MODEL_ID"] = model_id
    if model_url:
        os.environ["LANGEXTRACT_MODEL_URL"] = model_url


def _normalize_entity_name(name: str, mode: str = "aggressive") -> str:
    cleaned = name.strip().lower()
    if not cleaned:
        return ""
    cleaned = re.sub(r'^[\"`“”‘’\'()\[\]{}<>]+|[\"`“”‘’\'()\[\]{}<>]+$', "", cleaned)
    cleaned = cleaned.replace("_", " ").replace("/", " ").replace("-", " ")
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if mode == "aggressive":
        cleaned = re.sub(r"[^0-9a-zA-Z+#. ]+", " ", cleaned)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned


def _entity_id(
    name: str,
    ent_type: str,
    name_norm: str | None = None,
    project_id_normalized: str | None = None,
) -> str:
    """Generate a deterministic entity ID.

    Per Phase 04 of the unified ingest/query contract plan, the merge key
    is namespaced by ``project_id_normalized`` so entities from two projects
    that share one doc graph (e.g. ``{project_id}_doc``) stay distinct.
    ``project_id_normalized=None`` is preserved for backwards compatibility
    with single-project tests; it produces a globally-namespaced ID.
    """
    name_part = name_norm or _normalize_entity_name(name)
    if project_id_normalized:
        key = f"{project_id_normalized}::{ent_type}::{name_part}"
    else:
        key = f"{ent_type}::{name_part}"
    return str(uuid.uuid5(uuid.NAMESPACE_URL, key))


def build_graph_components(
    text: str,
    provider: str,
    nlp=None,
    gliner_model=None,
    gliner_labels: List[str] | None = None,
    gliner_threshold: float = 0.3,
    merge_entities: bool = True,
    normalize_mode: str = "aggressive",
    project_id_normalized: str | None = None,
) -> Tuple[Dict[str, Dict[str, str]], List[Dict[str, str]]]:
    if provider == "spacy":
        doc = nlp(text)
        entities = [{"name": ent.text, "type": ent.label_} for ent in doc.ents]
        relations = []
    elif provider == "gliner":
        entities, relations = extract_entities_gliner(
            text,
            labels=gliner_labels,
            threshold=gliner_threshold,
            gliner_model=gliner_model,
        )
    elif provider == "gemini":
        entities, relations = extract_entities_gemini(text)
    else:
        entities, relations = extract_entities_langextract(text)
    return build_graph_components_from_entities(
        entities,
        relations,
        merge_entities=merge_entities,
        normalize_mode=normalize_mode,
        project_id_normalized=project_id_normalized,
    )


def build_graph_components_from_entities(
    entities: List[Dict[str, str]],
    relations: List[Dict[str, str]],
    merge_entities: bool = True,
    normalize_mode: str = "aggressive",
    project_id_normalized: str | None = None,
) -> Tuple[Dict[str, Dict[str, str]], List[Dict[str, str]]]:
    nodes: Dict[str, Dict[str, str]] = {}
    name_index: Dict[str, str] = {}
    for ent in entities:
        name = ent.get("name", "").strip()
        if not name:
            continue
        ent_type = ent.get("type", "UNKNOWN") or "UNKNOWN"
        name_norm = _normalize_entity_name(name, mode=normalize_mode)
        if not name_norm:
            continue
        mention = {}
        if ent.get("start_char") is not None and ent.get("end_char") is not None:
            mention["start_char"] = ent.get("start_char")
            mention["end_char"] = ent.get("end_char")
            mention["surface"] = name
        if ent.get("confidence") is not None:
            mention["confidence"] = ent.get("confidence")
        if merge_entities:
            # Per-project merge key: ``{project_id_normalized}::{type}::{name}``.
            # When ``project_id_normalized`` is None we fall back to the legacy
            # ``{type}::{name}`` form so single-project tests keep passing.
            key_proj = f"{project_id_normalized}::{ent_type}::{name_norm}" if project_id_normalized else f"{ent_type}::{name_norm}"
            if key_proj not in nodes:
                nodes[key_proj] = {
                    "id": _entity_id(
                        name,
                        ent_type,
                        name_norm=name_norm,
                        project_id_normalized=project_id_normalized,
                    ),
                    "name": name,
                    "type": ent_type,
                    "name_norm": name_norm,
                    "mentions": [mention] if mention else [],
                    "confidence": mention.get("confidence") if mention else None,
                }
            else:
                if len(name) > len(nodes[key_proj]["name"]):
                    nodes[key_proj]["name"] = name
                if mention:
                    nodes[key_proj]["mentions"].append(mention)
                    conf = mention.get("confidence")
                    if conf is not None:
                        existing = nodes[key_proj].get("confidence")
                        if existing is None or conf > existing:
                            nodes[key_proj]["confidence"] = conf
            name_index.setdefault(name_norm, key_proj)
        else:
            key = str(uuid.uuid4())
            nodes[key] = {
                "id": str(uuid.uuid4()),
                "name": name,
                "type": ent_type,
                "name_norm": name_norm,
                "mentions": [mention] if mention else [],
                "confidence": mention.get("confidence") if mention else None,
            }
            name_index.setdefault(name_norm, key)
    cleaned_relations = []
    for rel in relations:
        source = rel.get("source", "").strip()
        target = rel.get("target", "").strip()
        rel_type = rel.get("relation", "").strip() or "RELATED"
        if not source or not target:
            continue
        source_norm = _normalize_entity_name(source, mode=normalize_mode)
        target_norm = _normalize_entity_name(target, mode=normalize_mode)
        if not source_norm or not target_norm:
            continue
        source_key = name_index.get(source_norm)
        if source_key is None:
            if merge_entities:
                source_key = (
                    f"{project_id_normalized}::UNKNOWN::{source_norm}"
                    if project_id_normalized
                    else f"UNKNOWN::{source_norm}"
                )
            else:
                source_key = str(uuid.uuid4())
            nodes[source_key] = {
                "id": _entity_id(
                    source,
                    "UNKNOWN",
                    name_norm=source_norm,
                    project_id_normalized=project_id_normalized,
                )
                if merge_entities
                else str(uuid.uuid4()),
                "name": source,
                "type": "UNKNOWN",
                "name_norm": source_norm,
                "mentions": [],
                "confidence": None,
            }
            name_index.setdefault(source_norm, source_key)
        target_key = name_index.get(target_norm)
        if target_key is None:
            if merge_entities:
                target_key = (
                    f"{project_id_normalized}::UNKNOWN::{target_norm}"
                    if project_id_normalized
                    else f"UNKNOWN::{target_norm}"
                )
            else:
                target_key = str(uuid.uuid4())
            nodes[target_key] = {
                "id": _entity_id(
                    target,
                    "UNKNOWN",
                    name_norm=target_norm,
                    project_id_normalized=project_id_normalized,
                )
                if merge_entities
                else str(uuid.uuid4()),
                "name": target,
                "type": "UNKNOWN",
                "name_norm": target_norm,
                "mentions": [],
                "confidence": None,
            }
            name_index.setdefault(target_norm, target_key)
        cleaned_relations.append(
            {
                "source_id": nodes[source_key]["id"],
                "target_id": nodes[target_key]["id"],
                "type": rel_type,
            }
        )
    return nodes, cleaned_relations


def _select_primary_mention(mentions: List[Dict[str, Any]]) -> Dict[str, Any]:
    if not mentions:
        return {}
    sorted_mentions = sorted(
        mentions,
        key=lambda m: (
            -(m.get("confidence") or 0.0),
            m.get("start_char") if m.get("start_char") is not None else 10**9,
        ),
    )
    return sorted_mentions[0]


def ingest_to_graph(
    driver,
    nodes: Dict[str, Dict[str, str]],
    relations: List[Dict[str, str]],
    source_id: str,
    paragraph_id: int,
    paragraph_text: str | None = None,
    is_short: bool = False,
    paragraph_props: Dict[str, Any] | None = None,
    project_id: str | None = None,
    project_id_normalized: str | None = None,
) -> None:
    with driver.session() as session:
        if paragraph_text:
            props = paragraph_props or {}
            session.run(
                """
                MERGE (d:Document {id: $doc_id})
                SET d.name               = $doc_name,
                    d.project_id         = $project_id,
                    d.project_id_normalized = $project_id_normalized
                MERGE (p:Paragraph {source_id: $source_id, paragraph_id: $paragraph_id})
                SET p.text               = $text,
                    p.short              = $is_short,
                    p.project_id         = $project_id,
                    p.project_id_normalized = $project_id_normalized
                SET p += $props
                MERGE (d)-[:HAS_PARAGRAPH]->(p)
                """,
                doc_id=source_id,
                doc_name=source_id,
                source_id=source_id,
                paragraph_id=paragraph_id,
                text=paragraph_text,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
                is_short=is_short,
                props=props,
            )
        for node in nodes.values():
            session.run(
                """
                MERGE (e:Entity {id: $id})
                SET e.name = coalesce(e.name, $name),
                    e.type = $type,
                    e.name_norm = $name_norm,
                    e.project_id = $project_id,
                    e.project_id_normalized = $project_id_normalized
                """,
                id=node["id"],
                name=node["name"],
                type=node["type"],
                name_norm=node["name_norm"],
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )
            if paragraph_text:
                primary = _select_primary_mention(node.get("mentions", []))
                session.run(
                    """
                    MATCH (p:Paragraph {source_id: $source_id, paragraph_id: $paragraph_id})
                    MATCH (e:Entity {id: $id})
                    MERGE (p)-[r:HAS_ENTITY]->(e)
                    SET r.source_id = $source_id,
                        r.paragraph_id = $paragraph_id,
                        r.confidence = $confidence,
                        r.start_char = $start_char,
                        r.end_char = $end_char,
                        r.surface = $surface,
                        r.project_id = $project_id,
                        r.project_id_normalized = $project_id_normalized
                    """,
                    source_id=source_id,
                    paragraph_id=paragraph_id,
                    id=node["id"],
                    confidence=node.get("confidence"),
                    start_char=primary.get("start_char"),
                    end_char=primary.get("end_char"),
                    surface=primary.get("surface"),
                    project_id=project_id,
                    project_id_normalized=project_id_normalized,
                )
        for rel in relations:
            session.run(
                """
                MATCH (s:Entity {id: $source_id})
                MATCH (t:Entity {id: $target_id})
                MERGE (s)-[r:RELATED {type: $type, source_id: $source_doc, paragraph_id: $paragraph_id}]->(t)
                SET r.project_id           = $project_id,
                    r.project_id_normalized = $project_id_normalized
                """,
                source_id=rel["source_id"],
                target_id=rel["target_id"],
                type=rel["type"],
                source_doc=source_id,
                paragraph_id=paragraph_id,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )


def ingest_to_graph_batch(
    driver,
    items: List[Dict[str, Any]],
    project_id: str | None = None,
    project_id_normalized: str | None = None,
) -> None:
    paragraphs = []
    entities = []
    relations = []
    for item in items:
        paragraph_text = item.get("paragraph_text")
        if paragraph_text is not None:
            paragraphs.append(
                {
                    "doc_id": item["source_id"],
                    "doc_name": item["source_id"],
                    "source_id": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                    "text": paragraph_text,
                    "is_short": item.get("is_short", False),
                    "props": item.get("paragraph_props") or {},
                }
            )
        for node in item.get("nodes", {}).values():
            primary = _select_primary_mention(node.get("mentions", []))
            entities.append(
                {
                    "id": node["id"],
                    "name": node["name"],
                    "type": node["type"],
                    "name_norm": node["name_norm"],
                    "source_id": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                    "confidence": node.get("confidence"),
                    "start_char": primary.get("start_char"),
                    "end_char": primary.get("end_char"),
                    "surface": primary.get("surface"),
                }
            )
        for rel in item.get("relations", []):
            relations.append(
                {
                    "source_id": rel["source_id"],
                    "target_id": rel["target_id"],
                    "type": rel["type"],
                    "source_doc": item["source_id"],
                    "paragraph_id": item["paragraph_id"],
                }
            )
    with driver.session() as session:
        if paragraphs:
            session.run(
                """
                UNWIND $paragraphs AS row
                MERGE (d:Document {id: row.doc_id})
                SET d.name                 = row.doc_name,
                    d.project_id           = $project_id,
                    d.project_id_normalized = $project_id_normalized
                MERGE (p:Paragraph {source_id: row.source_id, paragraph_id: row.paragraph_id})
                SET p.text                 = row.text,
                    p.short                = row.is_short,
                    p.project_id           = $project_id,
                    p.project_id_normalized = $project_id_normalized
                SET p += row.props
                MERGE (d)-[:HAS_PARAGRAPH]->(p)
                """,
                paragraphs=paragraphs,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )
        if entities:
            session.run(
                """
                UNWIND $entities AS row
                MERGE (e:Entity {id: row.id})
                SET e.name                 = coalesce(e.name, row.name),
                    e.type                 = row.type,
                    e.name_norm            = row.name_norm,
                    e.project_id           = $project_id,
                    e.project_id_normalized = $project_id_normalized
                """,
                entities=entities,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )
            session.run(
                """
                UNWIND $entities AS row
                MATCH (p:Paragraph {source_id: row.source_id, paragraph_id: row.paragraph_id})
                MATCH (e:Entity {id: row.id})
                MERGE (p)-[r:HAS_ENTITY]->(e)
                SET r.source_id            = row.source_id,
                    r.paragraph_id         = row.paragraph_id,
                    r.confidence           = row.confidence,
                    r.start_char           = row.start_char,
                    r.end_char             = row.end_char,
                    r.surface              = row.surface,
                    r.project_id           = $project_id,
                    r.project_id_normalized = $project_id_normalized
                """,
                entities=entities,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )
        if relations:
            session.run(
                """
                UNWIND $relations AS row
                MATCH (s:Entity {id: row.source_id})
                MATCH (t:Entity {id: row.target_id})
                MERGE (s)-[r:RELATED {type: row.type, source_id: row.source_doc, paragraph_id: row.paragraph_id}]->(t)
                SET r.project_id           = $project_id,
                    r.project_id_normalized = $project_id_normalized
                """,
                relations=relations,
                project_id=project_id,
                project_id_normalized=project_id_normalized,
            )


def create_collection(client: Any, name: str, vector_size: int) -> None:
    try:
        client.get_collection_info(name)
        return
    except Exception:
        pass
    client.create_collection(
        name,
        vectors_config=qmodels.VectorParams(size=vector_size, distance=qmodels.Distance.COSINE),
    )


_PRINTED_COLLECTIONS: set[str] = set()


def ingest_to_qdrant(
    client: Any,
    collection: str,
    paragraph: str,
    paragraph_id: int,
    embedder: SentenceTransformer,
    nodes: Dict[str, Dict[str, str]],
    source_id: str,
    project_id: str | None = None,
    project_id_normalized: str | None = None,
    extra_payload: Dict[str, Any] | None = None,
) -> None:
    vector = embedder.encode([paragraph])[0]
    entity_ids = list({node["id"] for node in nodes.values()})
    entity_mentions: List[Dict[str, Any]] = []
    for node in nodes.values():
        primary = _select_primary_mention(node.get("mentions", []))
        entity_mentions.append(
            {
                "id": node["id"],
                "name": node["name"],
                "type": node["type"],
                "start_char": primary.get("start_char"),
                "end_char": primary.get("end_char"),
                "confidence": node.get("confidence"),
            }
    )
    payload = {
        "paragraph_id": paragraph_id,
        "source_id": source_id,
        "text": paragraph,
        "entity_ids": entity_ids,
        "entity_mentions": entity_mentions,
    }
    if extra_payload:
        payload.update(extra_payload)
    payload["project_id"] = project_id
    payload["project_id_normalized"] = project_id_normalized
    point = qmodels.PointStruct(
        id=str(uuid.uuid4()),
        vector=vector.tolist(),
        payload=payload,
    )
    if collection not in _PRINTED_COLLECTIONS:
        print(f"Upserting into Qdrant collection: {collection}")
        _PRINTED_COLLECTIONS.add(collection)
    client.upsert(collection_name=collection, points=[point])


def _safe_source_id(base: Path, file_path: Path) -> str:
    rel = file_path.relative_to(base).as_posix()
    return rel.replace("/", "__").replace("\\", "__")


def _stringify_values(values: Dict[str, Any]) -> Dict[str, Any]:
    safe: Dict[str, Any] = {}
    for key, value in values.items():
        if value is None:
            safe[key] = None
        else:
            safe[key] = str(value)
    return safe


def _iter_input_files(folder: Path) -> List[Path]:
    files = []
    for path in folder.rglob("*"):
        if (
            path.is_file()
            and not path.name.startswith("~$")
            and path.suffix.lower() in {
                ".pdf",
                ".txt",
                ".md",
                ".docx",
                ".pptx",
                ".xlsx",
            }
        ):
            files.append(path)
    return sorted(files)


def _read_input_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf_text(file_path)
    if suffix == ".docx":
        return read_docx_text(file_path)
    if suffix == ".pptx":
        return read_pptx_text(file_path)
    if suffix == ".xlsx":
        return read_xlsx_text(file_path)
    return read_text_file(file_path)


def process_text(
    raw_text: str,
    source_id: str,
    args: argparse.Namespace,
    driver,
    qdrant: Any,
    embedder: SentenceTransformer,
    nlp=None,
    gliner_model=None,
    project_id_normalized: str | None = None,
) -> None:
    if not raw_text:
        print(f"Skip empty input: {source_id}")
        return

    ingested_at = datetime.now(timezone.utc).isoformat()
    print(f"Source: {source_id} at {ingested_at}")
    print(f"Entity provider: {args.entity_provider}")

    if args.entity_provider == "spacy" and nlp is None:
        nlp = build_spacy_pipeline(args.spacy_model, ruler_json=args.ruler_json)
    if args.entity_provider == "gliner" and gliner_model is None:
        gliner_model = build_gliner_model(args.gliner_model_resolved)
    gliner_labels = parse_gliner_labels(args.gliner_labels)
    use_gliner_batch = args.entity_provider == "gliner" and not args.no_batch
    merge_entities = not args.no_entity_merge
    normalize_mode = args.entity_normalize_mode
    if args.no_batch:
        args.gliner_batch_size = 1
        args.graph_batch_size = 1

    paragraphs = split_paragraphs(raw_text, args.max_paragraph_chars)
    print(f"Chunked into {len(paragraphs)} paragraphs.")
    graph_batch: List[Dict[str, Any]] = []

    def flush_graph_batch() -> None:
        if not graph_batch:
            return
        print(f"Ingesting {len(graph_batch)} paragraphs to graph store (batch)...")
        ingest_to_graph_batch(
            driver,
            graph_batch,
            project_id=(args.project_id or args.source_id or None),
            project_id_normalized=project_id_normalized,
        )
        graph_batch.clear()
        print("Graph batch ingestion complete.")

    long_paragraphs: List[Tuple[int, str]] = []
    for idx, paragraph in enumerate(paragraphs):
        if not paragraph:
            continue
        if len(paragraph) < args.min_paragraph_chars:
            if args.skip_llm_short:
                print(
                    f"Paragraph {idx + 1}/{len(paragraphs)}: LLM skipped (len={len(paragraph)})"
                )
                graph_batch.append(
                    {
                        "source_id": source_id,
                        "paragraph_id": idx,
                        "paragraph_text": paragraph,
                        "is_short": True,
                        "nodes": {},
                        "relations": [],
                    }
                )
                if len(graph_batch) >= args.graph_batch_size:
                    flush_graph_batch()
                print("Ingesting to Qdrant...")
                ingest_to_qdrant(
                    qdrant, args.collection, paragraph, idx, embedder, {}, source_id,
                    project_id=(args.project_id or args.source_id or None),
                    project_id_normalized=project_id_normalized,
                )
                print("Qdrant ingestion complete.")
            else:
                print(f"Paragraph {idx + 1}/{len(paragraphs)}: skipped (len={len(paragraph)})")
            continue
        long_paragraphs.append((idx, paragraph))

    if use_gliner_batch and long_paragraphs:
        batch_size = max(1, args.gliner_batch_size)
        for start in range(0, len(long_paragraphs), batch_size):
            batch_items = long_paragraphs[start : start + batch_size]
            batch_texts = [item[1] for item in batch_items]
            print(
                f"Extracting entities with GLiNER (batch {start + 1}-"
                f"{min(start + batch_size, len(long_paragraphs))})..."
            )
            batch_entities = extract_entities_gliner_batch(
                batch_texts,
                labels=gliner_labels,
                threshold=args.gliner_threshold,
                gliner_model=gliner_model,
            )
            for (idx, paragraph), entities in zip(batch_items, batch_entities, strict=True):
                nodes, relations = build_graph_components_from_entities(
                    entities,
                    [],
                    merge_entities=merge_entities,
                    normalize_mode=normalize_mode,
                    project_id_normalized=project_id_normalized,
                )
                print(f"Paragraph {idx + 1}/{len(paragraphs)}: extracted {len(nodes)} entities.")
                graph_batch.append(
                    {
                        "source_id": source_id,
                        "paragraph_id": idx,
                        "paragraph_text": paragraph,
                        "is_short": False,
                        "nodes": nodes,
                        "relations": relations,
                    }
                )
                if len(graph_batch) >= args.graph_batch_size:
                    flush_graph_batch()
                print("Ingesting to Qdrant...")
                ingest_to_qdrant(
                    qdrant, args.collection, paragraph, idx, embedder, nodes, source_id,
                    project_id=(args.project_id or args.source_id or None),
                    project_id_normalized=project_id_normalized,
                )
                print("Qdrant ingestion complete.")
    else:
        for idx, paragraph in long_paragraphs:
            print(f"Paragraph {idx + 1}/{len(paragraphs)}: extracting entities/relations...")
            nodes, relations = build_graph_components(
                paragraph,
                args.entity_provider,
                nlp=nlp,
                gliner_model=gliner_model,
                gliner_labels=gliner_labels,
                gliner_threshold=args.gliner_threshold,
                merge_entities=merge_entities,
                normalize_mode=normalize_mode,
                project_id_normalized=project_id_normalized,
            )
            print(f"Extracted {len(nodes)} entities, {len(relations)} relations.")
            graph_batch.append(
                {
                    "source_id": source_id,
                    "paragraph_id": idx,
                    "paragraph_text": paragraph,
                    "is_short": False,
                    "nodes": nodes,
                    "relations": relations,
                }
            )
            if len(graph_batch) >= args.graph_batch_size:
                flush_graph_batch()
            print("Ingesting to Qdrant...")
            ingest_to_qdrant(
                qdrant, args.collection, paragraph, idx, embedder, nodes, source_id,
                project_id=(args.project_id or args.source_id or None),
                project_id_normalized=project_id_normalized,
            )
            print("Qdrant ingestion complete.")
    flush_graph_batch()


def process_xlsx_structured(
    xlsx_path: Path,
    source_id: str,
    args: argparse.Namespace,
    driver,
    qdrant: Any,
    embedder: SentenceTransformer,
    project_id_normalized: str | None = None,
) -> None:
    from extractor.excel.excel_table_pipeline_skeleton import (
        apply_merged_cell_fill,
        build_cell_map,
        detect_header_rows,
        detect_tables,
        extract_entities_for_rows,
        extract_merged_ranges,
        extract_table_matrix,
        iter_data_rows,
        load_workbooks,
        normalize_headers,
        parse_entity_column_map,
        parse_entity_columns,
    )

    key_columns = [c.strip() for c in (args.xlsx_key_columns or "").split(",") if c.strip()]
    entity_columns = parse_entity_columns(args.xlsx_entity_columns)
    entity_column_map = parse_entity_column_map(args.xlsx_entity_column_map)
    if args.xlsx_no_entities:
        entity_provider = "none"
    else:
        entity_provider = args.xlsx_entity_provider or args.entity_provider
    entity_batch_size = args.xlsx_entity_batch_size or args.gliner_batch_size

    wb_values, wb_formula = load_workbooks(xlsx_path)
    sheets = [wb_values[args.xlsx_sheet]] if args.xlsx_sheet else list(wb_values.worksheets)

    graph_batch: List[Dict[str, Any]] = []

    def flush_graph_batch() -> None:
        if not graph_batch:
            return
        ingest_to_graph_batch(
            driver,
            graph_batch,
            project_id=(args.project_id or args.source_id or None),
            project_id_normalized=project_id_normalized,
        )
        graph_batch.clear()

    row_counter = 0
    for sheet in sheets:
        sheet_formula = wb_formula[sheet.title]
        cells = build_cell_map(sheet, sheet_formula)
        merged_ranges = extract_merged_ranges(sheet_formula)
        apply_merged_cell_fill(cells, merged_ranges)
        tables = detect_tables(
            cells,
            sheet.title,
            min_rows=args.xlsx_min_rows,
            min_cols=args.xlsx_min_cols,
            max_row_gap=args.xlsx_max_row_gap,
            max_col_gap=args.xlsx_max_col_gap,
        )
        print(f"Sheet: {sheet.title} -> tables: {len(tables)}")
        for table in tables:
            matrix = extract_table_matrix(cells, table)
            header_rows = detect_header_rows(
                matrix,
                max_header_rows=args.xlsx_max_header_rows,
                min_text_ratio=args.xlsx_min_text_ratio,
                max_numeric_ratio=args.xlsx_max_numeric_ratio,
            )
            headers = normalize_headers(matrix, header_rows)
            rows = list(
                iter_data_rows(
                    matrix,
                    headers,
                    header_rows,
                    table.table_id,
                    key_columns=key_columns,
                    skip_footer_rows=args.xlsx_skip_footer_rows,
                )
            )
            if not rows:
                continue
            if entity_provider != "none" or entity_columns or entity_column_map:
                entities_by_row = extract_entities_for_rows(
                    rows,
                    provider=entity_provider,
                    batch_size=entity_batch_size,
                    gliner_labels=args.gliner_labels,
                    gliner_threshold=args.gliner_threshold,
                    gliner_model=args.gliner_model_resolved,
                    spacy_model=args.spacy_model,
                    spacy_ruler=args.ruler_json,
                    column_map=entity_column_map,
                    column_list=entity_columns,
                    column_default_type=args.xlsx_entity_column_type,
                )
            else:
                entities_by_row = [[] for _ in rows]

            for row, entities in zip(rows, entities_by_row, strict=True):
                nodes, relations = build_graph_components_from_entities(
                    entities,
                    [],
                    merge_entities=not args.no_entity_merge,
                    normalize_mode=args.entity_normalize_mode,
                    project_id_normalized=project_id_normalized,
                )
                row_counter += 1
                paragraph_text = row.serialized
                paragraph_props = {
                    "sheet_name": sheet.title,
                    "table_id": table.table_id,
                    "row_id": row.row_id,
                    "row_hash": row.row_hash,
                    "row_index": row.row_index,
                    "is_footer": row.is_footer,
                }
                extra_payload = {
                    "sheet_name": sheet.title,
                    "table_id": table.table_id,
                    "row_id": row.row_id,
                    "row_hash": row.row_hash,
                    "row_index": row.row_index,
                    "is_footer": row.is_footer,
                    "values": row.values,
                    "raw_values": _stringify_values(row.raw_values),
                }
                ingest_to_qdrant(
                    qdrant,
                    args.collection,
                    paragraph_text,
                    row_counter,
                    embedder,
                    nodes,
                    source_id,
                    project_id=(args.project_id or args.source_id or None),
                    project_id_normalized=project_id_normalized,
                    extra_payload=extra_payload,
                )
                graph_batch.append(
                    {
                        "source_id": source_id,
                        "paragraph_id": row_counter,
                        "paragraph_text": paragraph_text,
                        "is_short": False,
                        "nodes": nodes,
                        "relations": relations,
                        "paragraph_props": paragraph_props,
                    }
                )
                if len(graph_batch) >= args.graph_batch_size:
                    flush_graph_batch()
    flush_graph_batch()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pdf", help="Path to PDF")
    parser.add_argument("--text-file", help="Path to UTF-8 text file (.txt)")
    parser.add_argument("--md", help="Path to UTF-8 Markdown file (.md)")
    parser.add_argument("--docx", help="Path to Word document (.docx)")
    parser.add_argument("--pptx", help="Path to PowerPoint file (.pptx)")
    parser.add_argument("--xlsx", help="Path to Excel file (.xlsx)")
    parser.add_argument(
        "--xlsx-structured",
        action="store_true",
        help="Use structured Excel ingestion (table/row aware).",
    )
    parser.add_argument("--xlsx-sheet", default=None, help="Sheet name for structured Excel")
    parser.add_argument(
        "--xlsx-key-columns",
        default="",
        help="Comma-separated key columns for stable row_id (structured Excel)",
    )
    parser.add_argument("--xlsx-min-rows", type=int, default=2)
    parser.add_argument("--xlsx-min-cols", type=int, default=2)
    parser.add_argument("--xlsx-max-row-gap", type=int, default=1)
    parser.add_argument("--xlsx-max-col-gap", type=int, default=1)
    parser.add_argument("--xlsx-max-header-rows", type=int, default=3)
    parser.add_argument("--xlsx-min-text-ratio", type=float, default=0.5)
    parser.add_argument("--xlsx-max-numeric-ratio", type=float, default=0.5)
    parser.add_argument("--xlsx-skip-footer-rows", action="store_true")
    parser.add_argument(
        "--xlsx-entity-provider",
        choices=["none", "spacy", "gliner", "gemini", "langextract"],
        default=None,
        help="Override entity provider for structured Excel (default: use --entity-provider).",
    )
    parser.add_argument(
        "--xlsx-entity-batch-size",
        type=int,
        default=None,
        help="Batch size for structured Excel entity extraction (gliner only).",
    )
    parser.add_argument(
        "--xlsx-entity-columns",
        default="",
        help="Comma-separated column names to treat as entities (structured Excel).",
    )
    parser.add_argument(
        "--xlsx-entity-column-map",
        default="",
        help="Comma-separated column:type mapping for entities (structured Excel).",
    )
    parser.add_argument(
        "--xlsx-entity-column-type",
        default="COLUMN",
        help="Default entity type for --xlsx-entity-columns.",
    )
    parser.add_argument(
        "--xlsx-no-entities",
        action="store_true",
        help="Disable entity extraction for structured Excel.",
    )
    parser.add_argument("--raw-text", help="Raw text input")
    parser.add_argument(
        "--folder",
        help="Folder to scan for .pdf/.txt/.md/.docx/.pptx/.xlsx files (recursive)",
    )
    parser.add_argument("--source-id", default=None, help="Custom source identifier")
    parser.add_argument(
        "--project-id",
        default=os.getenv("PROJECT_ID"),
        help=(
            "Project identifier. Stamped on Document/Paragraph/Entity nodes "
            "and on Qdrant payloads. When set, the ProjectRegistry resolves "
            "the doc graph and Qdrant collection unless an explicit "
            "--collection or --graph value is provided. When absent, "
            "fallback to source_id is used and a deprecation warning is "
            "logged."
        ),
    )
    parser.add_argument(
        "--collection",
        default=os.getenv("QDRANT_COLLECTION_DOC"),
    )
    parser.add_argument("--embedding-model", default=None)
    parser.add_argument("--embedding-device", default=None)
    parser.add_argument("--max-paragraph-chars", type=int, default=1200)
    parser.add_argument("--min-paragraph-chars", type=int, default=150)
    parser.add_argument(
        "--skip-llm-short",
        action="store_true",
        help="Skip LLM extraction for short paragraphs but still store in Qdrant",
    )
    parser.add_argument(
        "--entity-provider",
        choices=["spacy", "gliner", "gemini", "langextract"],
        default="gliner",
        help="Entity extraction provider",
    )
    parser.add_argument(
        "--entity-normalize-mode",
        choices=["basic", "aggressive"],
        default="aggressive",
        help="Normalization strength for entity merging.",
    )
    parser.add_argument("--spacy-model", default="en_core_web_sm")
    parser.add_argument("--ruler-json", default=None, help="Path to EntityRuler JSON patterns")
    parser.add_argument(
        "--gliner-model-name",
        default=os.getenv("GLINER_MODEL_NAME", "urchade/gliner_large-v2.1"),
        help="GLiNER model name (HuggingFace).",
    )
    parser.add_argument(
        "--gliner-model-path",
        default=os.getenv("GLINER_MODEL_PATH"),
        help="Local path to a GLiNER model directory.",
    )
    parser.add_argument(
        "--gliner-model",
        default=None,
        help="Deprecated: use --gliner-model-name or --gliner-model-path.",
    )
    parser.add_argument(
        "--gliner-labels",
        default="PERSON,ORG,PRODUCT,GPE,DATE,TECH,CRYPTO,STANDARD",
        help="Comma-separated labels or path to a text file for GLiNER",
    )
    parser.add_argument("--gliner-threshold", type=float, default=0.3)
    parser.add_argument(
        "--gliner-batch-size",
        type=int,
        default=8,
        help="Batch size for GLiNER entity extraction.",
    )
    parser.add_argument(
        "--no-entity-merge",
        action="store_true",
        help="Do not merge entities across paragraphs (use per-paragraph IDs).",
    )
    parser.add_argument(
        "--no-batch",
        action="store_true",
        help="Disable batching for GLiNER and graph-store (sequential processing).",
    )
    parser.add_argument(
        "--batch",
        action="store_false",
        dest="no_batch",
        help="Enable batching for GLiNER and graph-store.",
    )
    parser.add_argument("--langextract-model-id", default=os.getenv("LANGEXTRACT_MODEL_ID"))
    parser.add_argument("--langextract-model-url", default=os.getenv("LANGEXTRACT_MODEL_URL"))
    parser.add_argument("--llm-debug", action="store_true", help="Print raw LLM output")
    parser.add_argument("--llm-retry-count", type=int, default=None)
    parser.add_argument("--llm-retry-backoff", type=float, default=None)
    parser.add_argument(
        "--graph-batch-size",
        "--neo4j-batch-size",
        dest="graph_batch_size",
        type=int,
        default=50,
        help="Number of paragraphs to write per graph-store batch.",
    )
    add_graph_store_args(parser)
    parser.add_argument("--neo4j-uri", default=os.getenv("NEO4J_URI", "bolt://localhost:7687"))
    parser.add_argument("--neo4j-user", default=os.getenv("NEO4J_USER", "neo4j"))
    parser.add_argument("--neo4j-pass", default=os.getenv("NEO4J_PASS", "password"))
    parser.add_argument("--qdrant-path", default=os.getenv("QDRANT_DOC_PATH"))
    parser.set_defaults(no_batch=True)
    args = parser.parse_args()

    # Resolve ``project_id_normalized`` for the new contract. When the
    # caller passed --project-id we use it directly. When --project-id is
    # absent we fall back to --source-id (the legacy single-axis scope)
    # and log a deprecation warning so callers know to migrate.
    raw_project_id = (args.project_id or args.source_id or "").strip()
    if not args.project_id and args.source_id:
        print(
            "[graphrag_ingest_langextract] WARNING: --project-id not set; "
            "falling back to --source-id. The unified ingest/query contract "
            "requires --project-id. Pass --project-id explicitly to remove "
            "this warning and enable per-project entity isolation.",
            flush=True,
        )
    # Case-insensitive comparison key — mirrors tools.common.project_scope
    # in code-tiny. We inline the helper here because doc-tiny does not
    # have a Python package layout that can import from code-tiny.
    project_id_normalized = raw_project_id.casefold() if raw_project_id else None

    # Route ingest through the same registry contract used by mind_mcp. An
    # explicit CLI/env target remains an escape hatch; otherwise project_id
    # determines both the document graph and vector collection.
    if raw_project_id:
        try:
            targets = resolve_project_targets(raw_project_id)
        except ProjectNotRegisteredError as exc:
            raise SystemExit(str(exc)) from exc
        if not args.collection:
            args.collection = targets.doc_qdrant_collection
        graph_was_explicit = (
            "--falkordb-graph" in sys.argv
            or bool(os.getenv("FALKORDB_GRAPH"))
            or bool(os.getenv("FALKORDB_DATABASE"))
        )
        if not graph_was_explicit:
            args.falkordb_graph = targets.doc_graph
    if not args.collection:
        args.collection = (
            f"{raw_project_id}_doc" if raw_project_id else "graphrag_entities"
        )

    if args.gliner_model_path:
        gliner_model_choice = args.gliner_model_path
    elif args.gliner_model:
        gliner_model_choice = args.gliner_model
    else:
        gliner_model_choice = args.gliner_model_name
    args.gliner_model_resolved = gliner_model_choice

    inputs = [
        bool(args.pdf),
        bool(args.text_file),
        bool(args.md),
        bool(args.docx),
        bool(args.pptx),
        bool(args.xlsx),
        bool(args.raw_text),
        bool(args.folder),
    ]
    if sum(inputs) != 1:
        raise SystemExit(
            "Provide exactly one of --pdf, --text-file, --md, --docx, --pptx, --xlsx, --raw-text, or --folder."
        )

    if args.llm_debug:
        os.environ["LLM_DEBUG"] = "1"
        os.environ.setdefault("LLM_DEBUG_LOG_PATH", "logs/langextract_raw.log")
    if args.llm_retry_count is not None:
        os.environ["LLM_RETRY_COUNT"] = str(args.llm_retry_count)
    if args.llm_retry_backoff is not None:
        os.environ["LLM_RETRY_BACKOFF_SECONDS"] = str(args.llm_retry_backoff)
    _set_langextract_overrides(args.langextract_model_id, args.langextract_model_url)
    model_name, local_files_only = resolve_embedding_model(
        args.embedding_model, "BAAI/bge-m3"
    )
    device = resolve_embedding_device(args.embedding_device)
    embedder = SentenceTransformer(model_name, local_files_only=local_files_only, device=device)

    qdrant = get_document_qdrant_store(args.qdrant_path)
    create_collection(qdrant, args.collection, vector_size=embedder.get_sentence_embedding_dimension())

    driver = create_graph_store_from_args(args)

    if args.folder:
        folder_path = Path(args.folder)
        if not folder_path.exists():
            raise FileNotFoundError(folder_path)
        file_paths = _iter_input_files(folder_path)
        if not file_paths:
            raise SystemExit("No supported files found in folder.")
        print(f"Found {len(file_paths)} files in folder.")
        shared_nlp = None
        shared_gliner = None
        if args.entity_provider == "spacy":
            shared_nlp = build_spacy_pipeline(args.spacy_model, ruler_json=args.ruler_json)
        if args.entity_provider == "gliner":
            shared_gliner = build_gliner_model(args.gliner_model_resolved)
        for idx, file_path in enumerate(file_paths, start=1):
            print(f"[{idx}/{len(file_paths)}] Processing: {file_path}")
            if args.source_id:
                source_id = f"{args.source_id}__{_safe_source_id(folder_path, file_path)}"
            else:
                source_id = _safe_source_id(folder_path, file_path)
            raw_text = _read_input_text(file_path)
            process_text(
                raw_text,
                source_id,
                args,
                driver,
                qdrant,
                embedder,
                nlp=shared_nlp,
                gliner_model=shared_gliner,
                project_id_normalized=project_id_normalized,
            )
        return

    if args.pdf:
        pdf_path = Path(args.pdf)
        if not pdf_path.exists():
            raise FileNotFoundError(pdf_path)
        raw_text = read_pdf_text(pdf_path)
        source_id = args.source_id or pdf_path.stem
        process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return

    if args.text_file:
        text_path = Path(args.text_file)
        if not text_path.exists():
            raise FileNotFoundError(text_path)
        raw_text = read_text_file(text_path)
        source_id = args.source_id or text_path.stem
        process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return
    
    if args.md:
        md_path = Path(args.md)
        if not md_path.exists():
            raise FileNotFoundError(md_path)
        raw_text = read_text_file(md_path)
        source_id = args.source_id or md_path.stem
        process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return

    if args.docx:
        docx_path = Path(args.docx)
        if not docx_path.exists():
            raise FileNotFoundError(docx_path)
        raw_text = read_docx_text(docx_path)
        source_id = args.source_id or docx_path.stem
        process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return

    if args.pptx:
        pptx_path = Path(args.pptx)
        if not pptx_path.exists():
            raise FileNotFoundError(pptx_path)
        raw_text = read_pptx_text(pptx_path)
        source_id = args.source_id or pptx_path.stem
        process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return

    if args.xlsx:
        xlsx_path = Path(args.xlsx)
        if not xlsx_path.exists():
            raise FileNotFoundError(xlsx_path)
        source_id = args.source_id or xlsx_path.stem
        if args.xlsx_structured:
            process_xlsx_structured(
                xlsx_path,
                source_id,
                args,
                driver,
                qdrant,
                embedder,
                project_id_normalized=project_id_normalized,
            )
        else:
            raw_text = read_xlsx_text(xlsx_path)
            process_text(
            raw_text,
            source_id,
            args,
            driver,
            qdrant,
            embedder,
            project_id_normalized=project_id_normalized,
        )
        return

    raw_text = args.raw_text.strip()
    source_id = args.source_id or "raw_text"
    process_text(
        raw_text,
        source_id,
        args,
        driver,
        qdrant,
        embedder,
        project_id_normalized=project_id_normalized,
    )


if __name__ == "__main__":
    main()
